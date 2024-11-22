import os
from pptx import Presentation
from PIL import Image, ImageDraw

# Configuration
width, height = 1280, 720
image_path = "images"

# Convert PPT to Images
def ppt_to_images(ppt_path):
    prs = Presentation(ppt_path)
    slide_images = []
    if not os.path.exists(image_path):
        os.makedirs(image_path)
    
    for i, slide in enumerate(prs.slides):
        image_file = os.path.join(image_path, f"slide_{i + 1}.png")
        
        # Create a blank image for the slide
        img = Image.new('RGB', (width, height), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        y_offset = 50  # Vertical offset for text lines
        
        # Draw text from each shape in the slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Display text from the slide on the image
                    draw.text((50, y_offset), run.text, fill=(0, 0, 0))
                    y_offset += 40
        
        img.save(image_file)
        slide_images.append(image_file)
    return slide_images

def get_ppt_images(ppt_folder):
    ppt_files = [f for f in os.listdir(ppt_folder) if f.endswith(('.ppt', '.pptx'))]
    if not ppt_files:
        print("No PowerPoint files found in the folder.")
        return []
    ppt_path = os.path.join(ppt_folder, ppt_files[0])
    return ppt_to_images(ppt_path)
