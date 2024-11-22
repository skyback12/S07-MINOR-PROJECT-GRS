import cv2
import os
import numpy as np
from pptx import Presentation
from PIL import Image, ImageDraw
from HandTracker import HandDetector
import time

# Configuration
width, height = 1280, 720
ppt_folder = "PPT"
gesturethreshold = 500
buttondelay = 20  # Increased button delay to slow transitions
imgnumber = 0
buttonPressed, buttoncounter = False, 0
annotations = {}
annotationnumber, annotationstart = 0, False
image_path = "images"
zoom_factor = 1  # Default zoom factor

# Initialize Hand Detector
detector = HandDetector(detectionCon=0.8, maxHands=1)

# Convert PPT to Images
def ppt_to_images(ppt_path):
    prs = Presentation(ppt_path)
    slide_images = []
    if not os.path.exists(image_path):
        os.makedirs(image_path)
    
    for i, slide in enumerate(prs.slides):
        image_file = os.path.join(image_path, f"slide_{i + 1}.png")
        
        # Create a blank image for the slide
        img = Image.new('RGB', (width, height), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        y_offset = 50  # Vertical offset for text lines
        
        # Draw text from each shape in the slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Display text from the slide on the image
                    draw.text((50, y_offset), run.text, fill=(0, 0, 0))
                    y_offset += 40
        
        img.save(image_file)
        slide_images.append(image_file)
    return slide_images

# Load PowerPoint file and convert to images
ppt_files = [f for f in os.listdir(ppt_folder) if f.endswith(('.ppt', '.pptx'))]
if not ppt_files:
    print("No PowerPoint files found in the folder.")
else:
    ppt_path = os.path.join(ppt_folder, ppt_files[0])
    pathimages = ppt_to_images(ppt_path)
    print("Slides saved as images.")

# Setup Camera
cap = cv2.VideoCapture(0)
cap.set(3, width)
cap.set(4, height)

# Define the function to calculate the Euclidean distance between two landmarks
def calculate_distance(lmList, point1, point2):
    x1, y1, z1 = lmList[point1]
    x2, y2, z2 = lmList[point2]
    return np.sqrt((x2 - x1) ** 2 + (y2 - y1) ** 2)  # 2D distance

# Define thresholds for zoom gestures
min_distance = 30  # Adjust this threshold based on testing
max_distance = 150

while True:
    # Capture and flip the frame
    success, img = cap.read()
    img = cv2.flip(img, 1)

    # Resize image for display
    img = cv2.resize(img, (1240, 720))

    # Load and display the current slide image
    imgcurrent = cv2.imread(pathimages[imgnumber])
    if imgcurrent is None:
        print("Error loading slide image.")
        break

    # Apply zoom factor to the image
    imgcurrent = cv2.resize(imgcurrent, None, fx=zoom_factor, fy=zoom_factor)

    # Add a camera thumbnail
    hs, ws = int(120 * 1), int(213 * 1)
    imgSmall = cv2.resize(img, (ws, hs))
    h, w, _ = imgcurrent.shape
    imgcurrent[0:hs, w - ws:w] = imgSmall

    # Detect hands and process gestures
    hands, img = detector.findHands(img)
    if hands and not buttonPressed:
        hand = hands[0]
        fingers = detector.fingersUp(hand)
        cx, cy = hand['center']
        lmList = hand['lmList']
        xVal = int(np.interp(lmList[8][0], [width // 2, w], [0, width]))
        yVal = int(np.interp(lmList[8][1], [150, height - 150], [0, height]))
        indexFinger = xVal, yVal

        if cy <= gesturethreshold:
            # Previous Slide
            if fingers == [1, 0, 0, 0, 0] and imgnumber > 0:
                print("Left (Previous Slide)")
                buttonPressed = True
                annotations[imgnumber] = []
                annotationnumber = 0
                annotationstart = False
                imgnumber -= 1
                time.sleep(0.5)  # Add delay for smoother transition

            # Next Slide
            elif fingers == [0, 0, 0, 0, 1] and imgnumber < len(pathimages) - 1:
                print("Right (Next Slide)")
                buttonPressed = True
                annotations[imgnumber] = []
                annotationnumber = 0
                annotationstart = False
                imgnumber += 1
                time.sleep(0.5)  # Add delay for smoother transition

        # Pointer mode
        if fingers == [0, 1, 1, 0, 0]:
            cv2.circle(imgcurrent, indexFinger, 12, (0, 0, 255), cv2.FILLED)

        # Drawing mode
        if fingers == [0, 1, 0, 0, 0]:
            if not annotationstart:
                annotationstart = True
                annotationnumber += 1
                annotations[imgnumber] = []
            cv2.circle(imgcurrent, indexFinger, 12, (0, 255, 0), cv2.FILLED)
            annotations[imgnumber].append(indexFinger)
        else:
            annotationstart = False

        # Clear Screen
        if fingers == [1, 1, 1, 1, 1] and annotations:
            annotations[imgnumber] = []
            annotationnumber = 0
            print("Screen Cleared")

        # Undo Last Annotation
        if fingers == [0, 1, 1, 1, 0] and annotationnumber >= 0:
            if annotations[imgnumber]:
                annotations[imgnumber].pop(-1)
                annotationnumber -= 1
                buttonPressed = True
                print("Undo Last Annotation")

        # Zoom Gesture
        if fingers == [1, 1, 0, 0, 0]:  # Pinch gesture
            dist = calculate_distance(lmList, 4, 8)  # Distance between thumb (4) and index finger (8)
            if dist < min_distance:
                zoom_factor = 1.2  # Zoom in
            elif dist > max_distance:
                zoom_factor = 0.8  # Zoom out

        # Screenshot
        if fingers == [0, 0, 0, 0, 0]:  # Fist gesture to capture screenshot
            cv2.imwrite(f"screenshot_{imgnumber + 1}.png", imgcurrent)
            print(f"Screenshot of slide {imgnumber + 1} saved.")

    # Draw annotations on the slide
    for annotation in annotations.get(imgnumber, []):
        for i in range(1, len(annotation)):
            cv2.line(imgcurrent, annotation[i - 1], annotation[i], (0, 0, 200), 12)

    # Display slide number
    cv2.putText(imgcurrent, f"Slide {imgnumber + 1}/{len(pathimages)}", (w - 150, h - 20),
                cv2.FONT_HERSHEY_SIMPLEX, 1, (0, 255, 0), 2)

    # Display images
    cv2.imshow("Slides", imgcurrent)
    key = cv2.waitKey(1)
    if key == ord('q'):
        break

    # Reset button press state
    if buttonPressed:
        buttoncounter -= 1
        if buttoncounter <= 0:
            buttonPressed = False

    # Add a delay to slow down the loop
    time.sleep(0.1)

cap.release()
cv2.destroyAllWindows()
